#!/usr/bin/env python3
"""Export GMMT course posters (23 PNGs) into a single 16:9 PPTX."""
import os
from pptx import Presentation
from pptx.util import Inches

POSTER_DIR = "/Users/hongyangchun/Codebase/default/outputs/gmmt-ppt/frontend/public/assets/images/posters/pages"
OUT = "/Users/hongyangchun/Codebase/default/outputs/gmmt-ppt/GLMT模型课程.pptx"
TOTAL = 23

prs = Presentation()
# 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

added = 0
for i in range(1, TOTAL + 1):
    path = os.path.join(POSTER_DIR, f"page-{i}.png")
    if not os.path.exists(path):
        print(f"[WARN] missing: {path}")
        continue
    slide = prs.slides.add_slide(blank)
    # stretch poster to fill the whole slide
    slide.shapes.add_picture(path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    added += 1

prs.save(OUT)
print(f"[OK] saved {OUT}  slides={added}")
